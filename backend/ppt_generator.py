# backend/ppt_generator.py

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

def create_pptx_from_slides(slides: list, output_path: str):
    """
    Creates a PowerPoint presentation from slide content.

    Args:
        slides (list): List of dicts like {"title": "...", "bullets": ["...", "..."]}
        output_path (str): Path to save the .pptx file
    """
    prs = Presentation()

    # Set a basic slide layout: Title + Content
    layout = prs.slide_layouts[1]

    for slide in slides:
        title = slide.get("title", "Untitled Slide")
        bullets = slide.get("bullets", [])

        # Create a new slide
        slide_layout = prs.slides.add_slide(layout)
        title_placeholder = slide_layout.shapes.title
        content_placeholder = slide_layout.placeholders[1]

        # Set slide title
        title_placeholder.text = title

        # Add bullet points
        tf = content_placeholder.text_frame
        tf.clear()  # remove default bullet

        for i, bullet in enumerate(bullets):
            if i == 0:
                tf.text = bullet  # First bullet
            else:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

            # Optional styling
            for paragraph in tf.paragraphs:
                paragraph.font.size = Pt(18)

    # Save the presentation
    prs.save(output_path)
